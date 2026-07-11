import os
import pdfplumber
from docx import Document
from pptx import Presentation


def extract_from_pdf(path: str) -> list[dict]:
    """Extrait le texte page par page d'un PDF."""
    units = []
    filename = os.path.basename(path)

    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                units.append({
                    "filename": filename,
                    "unit_type": "page",
                    "unit_number": i,
                    "text": text,
                })
    return units


def extract_from_docx(path: str) -> list[dict]:
    """Extrait le texte d'un Word, regroupé par blocs de 10 paragraphes."""
    filename = os.path.basename(path)
    doc = Document(path)

    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    units = []
    block_size = 10
    for i in range(0, len(paragraphs), block_size):
        block_text = "\n".join(paragraphs[i:i + block_size])
        units.append({
            "filename": filename,
            "unit_type": "paragraph_block",
            "unit_number": (i // block_size) + 1,
            "text": block_text,
        })
    return units


def extract_from_pptx(path: str) -> list[dict]:
    """Extrait le texte slide par slide d'un PowerPoint."""
    filename = os.path.basename(path)
    prs = Presentation(path)

    units = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        texts.append(line.strip())

        slide_text = "\n".join(texts)
        if slide_text:
            units.append({
                "filename": filename,
                "unit_type": "slide",
                "unit_number": i,
                "text": slide_text,
            })
    return units


def extract_from_txt(path: str) -> list[dict]:
    """Extrait un fichier texte brut en une seule unité."""
    filename = os.path.basename(path)
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read().strip()

    if not text:
        return []

    return [{
        "filename": filename,
        "unit_type": "document",
        "unit_number": 1,
        "text": text,
    }]


def extract_document(path: str) -> list[dict]:
    """Point d'entrée unique : détecte l'extension et appelle la bonne fonction."""
    ext = path.split(".")[-1].lower()

    if ext == "pdf":
        return extract_from_pdf(path)
    elif ext in ("docx", "doc"):
        return extract_from_docx(path)
    elif ext in ("pptx", "ppt"):
        return extract_from_pptx(path)
    elif ext == "txt":
        return extract_from_txt(path)
    else:
        raise ValueError(f"Format non supporté : .{ext}")